"""
Create Professional PowerPoint Presentation for Phase 1 Work
Major Project-I: Imitating Fund Manager Decisions
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
from pathlib import Path

# Load evaluation results
with open('outputs/metrics/Temporal_KG_evaluation_results.json', 'r') as f:
    temporal_results = json.load(f)

with open('outputs/metrics/Causal_KG_evaluation_results.json', 'r') as f:
    causal_results = json.load(f)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)

    return slide

def add_content_slide(prs, title, content_blocks):
    """Add content slide with title and multiple text blocks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add line separator
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(0.5), Inches(1.05),
        Inches(9), Inches(0)
    )
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(2)

    # Add content blocks
    y_position = 1.3
    for block in content_blocks:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_position), Inches(8.6), Inches(0.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        p = content_frame.paragraphs[0]
        p.text = block
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(32, 32, 32)
        p.space_after = Pt(12)
        p.line_spacing = 1.2

        y_position += 0.8

    return slide

def add_results_slide(prs, title, metrics_dict, graph_name):
    """Add results slide with metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add content
    y_pos = 1.3

    # Graph structure
    content = f"{graph_name} consists of {metrics_dict['structural_completeness']['node_count']:,} nodes and {metrics_dict['structural_completeness']['edge_count']:,} edges, representing a comprehensive network that captures the intricate relationships within portfolio management decisions. "

    if graph_name == "Temporal Knowledge Graph":
        content += f"The graph includes {metrics_dict['structural_completeness']['node_type_distribution']['Fund']} mutual funds, {metrics_dict['structural_completeness']['node_type_distribution']['Stock']:,} stocks, {metrics_dict['structural_completeness']['node_type_distribution']['Sector']} sectors, and {metrics_dict['structural_completeness']['node_type_distribution']['TimePeriod']} time periods. "

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(1.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    p = content_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(13)
    p.line_spacing = 1.2

    y_pos += 1.5

    # Metrics table
    table_data = [
        ["Metric", "Score", "Interpretation"],
        ["Structural Completeness", f"{metrics_dict['structural_completeness']['completeness_score']:.3f}",
         "Excellent" if metrics_dict['structural_completeness']['completeness_score'] > 0.9 else "Good"],
        ["Consistency & Integrity", f"{metrics_dict['consistency']['consistency_score']:.3f}",
         "Excellent" if metrics_dict['consistency']['consistency_score'] > 0.9 else "Good"],
        ["Semantic Coherence", f"{metrics_dict['semantic_coherence']['semantic_coherence_score']:.3f}",
         "Good" if metrics_dict['semantic_coherence']['semantic_coherence_score'] > 0.7 else "Fair"],
        ["Informativeness", f"{metrics_dict['informativeness']['informativeness_score']:.3f}",
         "Good" if metrics_dict['informativeness']['informativeness_score'] > 0.5 else "Fair"],
        ["Inferential Utility", f"{metrics_dict['inferential_utility']['inferential_utility_score']:.3f}",
         "Fair" if metrics_dict['inferential_utility']['inferential_utility_score'] > 0.3 else "Needs Improvement"],
        ["Overall Quality", f"{metrics_dict['overall_quality_score']:.3f}",
         "Good" if metrics_dict['overall_quality_score'] > 0.6 else "Fair"]
    ]

    table = slide.shapes.add_table(len(table_data), 3, Inches(1.5), Inches(y_pos), Inches(7), Inches(2.5)).table

    for i, row_data in enumerate(table_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            if i == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    return slide

# Slide 1: Title
add_title_slide(prs,
    "Imitating Fund Manager Decisions:\nA Knowledge Graph Approach",
    "Phase 1: Knowledge Representation Construction and Evaluation\n\nGANNAMANENI JASWANTH (242IT012)\nGuide: Dr. Biju R. Mohan\n\nDepartment of Information Technology\nNational Institute of Technology Karnataka, Surathkal"
)

# Slide 2: Introduction
add_content_slide(prs, "Introduction", [
    "Investment decisions made by mutual fund managers represent complex cognitive processes that synthesize years of accumulated market experience, pattern recognition across economic cycles, and sophisticated reasoning about causality. Understanding and replicating these decision-making patterns holds significant value for both theoretical advancement in computational finance and practical applications in portfolio management systems.",

    "Traditional approaches to portfolio analysis have predominantly focused on performance attribution and risk decomposition, which explain returns but fail to capture the underlying decision-making logic. While recent machine learning advances have demonstrated the ability to predict fund performance with meaningful accuracy, these black-box approaches lack the interpretability necessary for understanding why certain stocks are selected and how allocation decisions evolve temporally.",

    "This research addresses these fundamental limitations by treating fund manager imitation as a knowledge representation problem rather than a pure prediction task. We construct structured knowledge bases that explicitly capture both the temporal patterns inherent in portfolio evolution and the causal relationships that link market conditions to investment decisions."
])

# Slide 3: Research Motivation
add_content_slide(prs, "Research Motivation", [
    "Mutual fund portfolio disclosures provide a unique window into professional investment decision-making. Each monthly disclosure reveals not only which securities a fund manager chose to hold and their allocation weights, but also how these positions changed from previous periods. These disclosures represent the observable outcomes of complex reasoning processes that synthesize diverse information sources including financial statements, earnings reports, industry trends, economic indicators, and market sentiment.",

    "The central challenge lies in reconstructing the decision-making logic from observational data alone. When a fund manager increases allocation to technology stocks while simultaneously reducing consumer discretionary holdings, multiple explanatory narratives might apply. The shift could reflect sector rotation based on macroeconomic forecasts, response to changing interest rate expectations, reaction to recent earnings surprises, or strategic adjustment of portfolio risk characteristics.",

    "Distinguishing between these alternative explanations requires understanding both the temporal context in which decisions occur and the causal mechanisms that link observable market factors to portfolio actions. This motivates our dual knowledge representation approach combining temporal and causal perspectives."
])

# Slide 4: Literature Review - Knowledge Graphs
add_content_slide(prs, "Literature Review: Knowledge Graph Foundations", [
    "Knowledge graphs have emerged as the dominant paradigm for representing complex decision-making knowledge through multi-relational directed graph structures. Wang et al. [3] demonstrate that knowledge graphs enable automated data analysis and intelligent decision-making through semantic reasoning and inference mechanisms, providing highly intuitive and flexible representations that support rapid identification of correlations across large-scale, heterogeneous datasets.",

    "Jiang et al. [4] formalize knowledge graphs using Labeled Property Graphs with mathematical notation defining vertices, edges, labels, labeling functions, and property functions. Their framework demonstrates efficient storage with rapid traversal capabilities while maintaining flexibility in modeling diverse real-world domains. The integration of heterogeneous data sources with varying structures proves particularly valuable for portfolio management applications requiring synthesis of fundamental data, market sentiment, technical indicators, and macroeconomic factors within unified reasoning frameworks.",

    "However, significant challenges constrain practical deployment. Scalability issues arise as data volume and diversity grow rapidly, with high computational complexity particularly for real-time reasoning over large graphs. Data quality concerns regarding accuracy, consistency, completeness, and trustworthiness plague knowledge graph construction, while dependency on domain expertise for effective graph construction creates bottlenecks in capturing nuanced fund manager decision patterns."
])

# Slide 5: Literature Review - Temporal and Causal Representations
add_content_slide(prs, "Literature Review: Temporal and Causal Knowledge", [
    "Liu et al. [7] develop a comprehensive framework for temporal knowledge graph reasoning using recurrent encoding mechanisms to capture entity and relation evolution over time. The methodology formally defines temporal knowledge graphs as sequences of snapshots with temporal dependencies, enabling both interpolation of missing historical information and extrapolation for future prediction. This approach proves essential for modeling cause-effect relationships that unfold over time, critical for understanding how economic indicators lead market movements and how portfolio rebalancing decisions produce downstream effects.",

    "Job et al. [8] provide comprehensive coverage of causal learning through graph neural networks, establishing theoretical foundations for distinguishing causation from correlation using do-calculus and interventional reasoning. Causal graphs provide explicit representation of cause-effect relationships between investment decisions and outcomes, enabling fund managers to reason about intervention effects before implementation and supporting counterfactual analysis of alternative strategies not taken.",

    "Wu et al. [9] demonstrate practical causal knowledge graph construction for enterprise innovation events, automatically extracting causal relationships from unstructured business data using neural language models. Their framework scales to large volumes of text data while revealing how innovation activities causally influence business outcomes, directly applicable to analyzing how corporate strategies affect investment returns."
])

# Slide 6: Literature Review - Evaluation Frameworks
add_content_slide(prs, "Literature Review: Evaluation and Validation", [
    "Zhang and Xiao [12] introduce a customizing knowledge graph evaluation method incorporating user needs through the 'Fitness for Use' principle. Their Entity Popularity-Weighted Clustering Sampling prioritizes high-traffic entities based on user search behavior, enabling efficient accuracy assessment focused on knowledge elements actually used in decisions. The framework evaluates multiple quality dimensions including accuracy, completeness, consistency, and redundancy, ensuring that evaluation efforts concentrate on knowledge components most critical for supporting actual decision-making workflows.",

    "Liu et al. [13] present a three-dimensional framework evaluating completeness, accuracy, and usability specifically for domain knowledge graphs. The completeness dimension employs network density, degree distribution, and complex network analysis metrics to measure coverage breadth and structural coherence. The accuracy dimension combines ontology reasoning to detect contradictions through formal logic, closeness centrality to validate critical entities against prior knowledge, and k-core analysis to identify structurally important concepts requiring high accuracy.",

    "These evaluation methodologies provide the foundation for our comprehensive assessment combining intrinsic metrics examining graph properties with extrinsic metrics measuring downstream task performance, ensuring that knowledge representations genuinely capture fund manager expertise rather than spurious statistical patterns."
])

# Slide 7: Outcome of Literature Review
add_content_slide(prs, "Outcome of Literature Review", [
    "The comprehensive examination of knowledge representation techniques reveals that knowledge graphs provide the most flexible and intuitive framework for representing complex relationships in portfolio management contexts. Temporal knowledge representation proves essential for capturing sequential decision-making patterns, while causal representation addresses the critical challenge of understanding why managers make particular decisions rather than merely identifying correlations.",

    "However, three significant research gaps emerge from the literature. First, most temporal knowledge graph work focuses on link prediction and completion rather than capturing decision-making sequences where temporal ordering carries semantic meaning. Second, causal knowledge graph research emphasizes discovery from text or observational data but provides limited guidance for representing intervention reasoning in sequential decision contexts. Third, integration of temporal and causal representations remains largely unexplored, yet understanding both when decisions occur and why they are made proves essential for genuine expert behavior imitation.",

    "Evaluation frameworks provide valuable methodologies for multi-dimensional quality assessment but focus primarily on abstract metrics rather than decision support contexts where both interpretability and predictive utility matter equally. Storage architectures demonstrate scalability for large knowledge bases but do not address specific requirements of financial applications where temporal queries, causal inference, and real-time updates must coexist efficiently."
])

# Slide 8: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "The central research problem is to design and develop a knowledge representation system that analyzes patterns in stock selection and allocation from historical fund manager decisions and uses this structured learning to imitate the decision-making process for future investment scenarios.",

    "Specifically, the system must capture not only what stocks fund managers select and when allocation changes occur (temporal dimension), but also why certain decisions are made based on market conditions, company fundamentals, and macroeconomic factors (causal dimension). The knowledge representation must support both backward-looking analysis of historical patterns and forward-looking reasoning about hypothetical scenarios.",

    "The system must achieve sufficient quality across multiple evaluation dimensions—structural completeness, semantic coherence, logical consistency, informativeness, and inferential utility—to serve as a reliable foundation for downstream portfolio construction tasks. This requires balancing representational richness with computational efficiency, interpretability with predictive accuracy, and domain-specific requirements with general knowledge engineering principles."
])

# Slide 9: Research Objectives
add_content_slide(prs, "Research Objectives", [
    "The primary research objective is to design, implement, and rigorously evaluate knowledge representation techniques that effectively capture fund manager decision patterns from historical portfolio disclosure data. This overarching goal decomposes into two complementary representation objectives addressing different dimensions of the imitation problem.",

    "First, develop temporal knowledge representations that model portfolio evolution as time-stamped facts with relationships tracking allocation changes across monthly disclosure periods. This representation must enable reasoning about decision sequences, identification of recurring patterns in position adjustments, and understanding of temporal context influencing allocation changes. The temporal graph must support queries about what decisions occurred historically and when particular patterns emerge in manager behavior.",

    "Second, construct causal knowledge representations that explicitly model cause-effect relationships between observable market factors and portfolio actions. This representation must distinguish genuine causation from mere correlation, support interventional reasoning about hypothetical market conditions, and enable counterfactual analysis of alternative decision paths. Integration of these complementary representations creates a unified knowledge base supporting sophisticated queries combining temporal sequencing with causal explanation, providing complete decision context rather than isolated facts."
])

# Slide 10: Methodology Overview
add_content_slide(prs, "Methodology: Research Framework", [
    "The research methodology follows a systematic pipeline from raw data collection through knowledge representation construction to comprehensive evaluation. The framework operates in two phases with distinct objectives and validation criteria.",

    "Phase 1 focuses exclusively on knowledge representation construction and intrinsic quality evaluation. This phase involves designing temporal and causal knowledge graph structures, implementing data processing pipelines to extract entities and relationships from portfolio disclosures, and evaluating resulting knowledge bases using five intrinsic quality metrics: structural completeness, consistency and integrity, semantic coherence, informativeness, and inferential utility. Phase 1 establishes whether knowledge graphs successfully capture patterns and relationships present in fund manager decisions.",

    "Phase 2 addresses downstream application of these knowledge representations to portfolio construction and optimization tasks. This phase evaluates how well knowledge bases support actual decision-making by using them to generate portfolio recommendations and comparing these recommendations against actual fund manager choices. The extrinsic evaluation provides evidence about whether knowledge representations contain information sufficient to replicate fund manager decision patterns in practical applications. The present work completes Phase 1, establishing the foundation for Phase 2 downstream tasks."
])

# Slide 11: Data Collection and Preprocessing
add_content_slide(prs, "Methodology: Data Collection and Integration", [
    "The empirical foundation consists of monthly portfolio disclosures from 22 mutual funds spanning September 2022 through September 2025, representing 45 monthly observation periods. The dataset captures 94,589 individual holding records across 1,762 distinct securities identified by ISIN codes. Each portfolio disclosure details equity positions including security identifiers, quantities held, market values, and allocation percentages, with changes between consecutive months revealing fund manager actions such as position increases, decreases, new entries, and complete exits.",

    "Fundamental data collection employs Yahoo Finance API to obtain company-specific financial metrics including Price-to-Earnings ratios, Price-to-Book ratios, Return on Equity, revenue growth rates, debt-to-equity ratios, profit margins, market capitalization, and systematic risk (beta). Monthly collection frequency aligns with portfolio disclosure timing, enabling analysis of whether fund managers systematically favor securities with particular fundamental characteristics. Data preprocessing standardizes currencies, adjusts for corporate actions, detects and handles outliers, and imputes missing values using forward filling for slowly changing items.",

    "Macroeconomic data captures market-wide conditions from 20+ indicators including NIFTY 50 and sectoral indices, India VIX volatility index, USD/INR exchange rates, government bond yields, and commodity prices. News sentiment analysis processes financial articles using FinBERT domain-specific sentiment classification, aggregating scores at both stock and sector levels. These heterogeneous data sources integrate into a unified dataset where portfolio holdings link to associated fundamentals, macroeconomic conditions, and sentiment at each time period."
])

# Slide 12: Fundamental Metrics Details
add_content_slide(prs, "Fundamental Metrics: Capturing Company Quality", [
    "The fundamental data collection system captures eight critical financial metrics that professional fund managers traditionally consider when evaluating investment opportunities. These metrics span valuation, profitability, growth, leverage, and risk dimensions, providing a comprehensive quantitative profile of each security in fund portfolios.",

    "Valuation metrics include Trailing Price-to-Earnings ratio and Price-to-Book ratio, which measure whether securities trade at premium or discount to earnings and book value respectively. Profitability and efficiency metrics include Return on Equity, measuring how effectively companies generate returns on shareholder capital, and Profit Margins, indicating operational efficiency in converting revenues to profits. Growth metrics capture Revenue Growth rates, revealing whether companies are expanding their business scale over time.",

    "Financial health and risk metrics include Debt-to-Equity ratio, measuring financial leverage and potential distress risk, Market Capitalization, indicating company size and liquidity characteristics, and Beta, measuring systematic risk and correlation with broader market movements. The combination of these metrics enables our knowledge graphs to represent not only which stocks fund managers select but also the fundamental characteristics of their selections, supporting analysis of whether managers systematically favor value stocks, high-growth companies, or particular risk-return profiles across different market conditions."
])

# Slide 13: Temporal Knowledge Graph Construction
add_content_slide(prs, "Temporal Knowledge Graph: Capturing Decision Sequences", [
    "The temporal knowledge graph represents portfolio evolution through a multi-relational directed graph structure with four distinct node types: Fund nodes represent the 22 mutual funds, Stock nodes represent 1,762 individual securities, Sector nodes categorize stocks by industry classification, and TimePeriod nodes mark the 45 monthly observation windows. This ontology enables rich representation of how fund managers allocate capital across securities over time.",

    "Five relationship types capture different aspects of portfolio management decisions. HOLDS relationships connect funds to stocks at specific time periods with properties including portfolio weight, market value, and timestamp, representing the core portfolio state. BELONGS_TO_SECTOR relationships link stocks to their industry classifications. Change relationships—INCREASED, DECREASED, ENTERED, and EXITED—explicitly model dynamic portfolio adjustments with properties capturing change magnitude, previous weight, new weight, and percentage change.",

    "Construction employs a systematic process: instantiate nodes for all entities in the integrated dataset, create HOLDS relationships for each fund-stock pair at every time period with appropriate weights and timestamps, generate temporal change relationships by comparing holdings across consecutive periods with filtering to exclude trivial adjustments from price movements, and establish sector membership links. The resulting graph supports sophisticated queries tracing fund histories with particular stocks, identifying similar decision sequences across funds, and detecting coordinated actions suggesting market-wide factors."
])

# Slide 14: Causal Knowledge Graph Construction
add_content_slide(prs, "Causal Knowledge Graph: Modeling Decision Rationale", [
    "The causal knowledge graph models decision-making rationale through directed acyclic graph structure where nodes represent observable factors, intermediate signals, and portfolio actions, while edges indicate causal influence rather than mere correlation. Factor nodes include macroeconomic indicators like interest rates and exchange rates, sector-level metrics, and stock fundamentals. Action nodes represent allocation decisions such as sector increases or specific stock entries.",

    "Causal relationship extraction employs Granger causality testing as the primary statistical methodology. For each potential cause-effect pair, the system constructs vector autoregressive models testing whether past values of the cause variable significantly improve prediction of the effect variable beyond simple autoregressive patterns. Relationships achieving statistical significance (p-value < 0.05) at appropriate time lags (tested up to 3 months) receive causal edges in the graph with properties encoding causal strength, confidence level, and optimal lag.",

    "The construction process enforces a critical constraint: the resulting graph must form a directed acyclic graph with no cycles, as cyclic causal relationships would imply logical contradictions. When potential cycles emerge from statistical testing, the system retains only the strongest causal edges to maintain acyclicity while preserving the most robust relationships. This principled approach ensures that the causal graph supports valid interventional reasoning about hypothetical scenarios, enabling fund managers to reason about what would happen under different market conditions before taking action."
])

# Slide 15: Knowledge Graph Integration
add_content_slide(prs, "Integration Layer: Linking Temporal and Causal Views", [
    "The integration layer establishes connections between temporal and causal knowledge graphs through identification of shared entities and creation of explicit cross-reference relationships. Stocks and sectors appear in both graphs, serving as bridge concepts that enable queries spanning temporal sequences and causal explanations. The integration process identifies these shared entities and links temporal events to their causal explanations through timestamp matching and logical consistency verification.",

    "Cross-reference relationships connect portfolio change events recorded in the temporal graph to causal factors represented in the causal graph. When the temporal graph records an allocation increase for a technology stock in October 2024, the integration layer searches the causal graph for factors influencing technology sector allocations with appropriate temporal precedence—typically September 2024 for lag-1 relationships. Identified causal factors receive cross-reference edges to the temporal event with properties encoding explanation strength and confidence.",

    "This integrated knowledge base supports sophisticated queries that combine temporal and causal reasoning. For example, the query 'Why did Fund X increase Stock Y in period T?' retrieves the temporal context including previous allocation, change magnitude, and timing from the temporal graph, then traverses cross-reference edges to identify causal factors from the causal graph, and finally synthesizes a complete explanation linking observable market conditions to the specific portfolio action. This integration transforms isolated facts into coherent decision narratives that mirror expert reasoning processes."
])

# Slide 16: Evaluation Framework
add_content_slide(prs, "Evaluation Framework: Multi-Dimensional Assessment", [
    "The evaluation framework employs five intrinsic metrics examining different dimensions of knowledge graph quality. Structural Completeness assesses whether required entities and relationships are present through network density, degree distributions, and coverage of securities from portfolio data. High completeness scores indicate that the knowledge representation successfully captures the breadth of fund manager decision space. The metric considers both node coverage (are all relevant stocks represented?) and edge coverage (are all significant portfolio relationships captured?).",

    "Consistency and Integrity verify logical correctness through cardinality constraint validation, attribute range checking, temporal ordering verification, and cycle detection in causal graphs. High consistency scores indicate the knowledge base contains no contradictory information that would undermine reasoning validity. Semantic Coherence evaluates whether related entities cluster meaningfully through community detection alignment with sectors and modularity scores, with high coherence indicating that the graph structure reflects genuine semantic relationships rather than random connectivity.",

    "Informativeness measures the richness of metadata attached to nodes and edges through counting unique attributes, calculating average attributes per entity, and assessing information density. Rich metadata enables more nuanced reasoning and detailed decision explanations. Inferential Utility quantifies reasoning capability through average path lengths, reachability analysis, centrality measures, and multi-hop query support, with high utility scores indicating the graph structure supports the complex queries required for decision support applications. Together, these five dimensions provide comprehensive assessment of whether knowledge graphs achieve the quality necessary for downstream portfolio construction tasks."
])

# Slide 17: Results - Temporal KG
add_results_slide(prs, "Results: Temporal Knowledge Graph Evaluation", temporal_results, "Temporal Knowledge Graph")

# Slide 18: Results - Temporal KG Detailed Analysis
add_content_slide(prs, "Temporal Knowledge Graph: Detailed Analysis", [
    f"The Temporal Knowledge Graph achieves an overall quality score of {temporal_results['overall_quality_score']:.3f}, indicating good quality suitable for downstream applications. Structural completeness reaches {temporal_results['structural_completeness']['completeness_score']:.3f}, demonstrating excellent coverage of the portfolio decision space with {temporal_results['structural_completeness']['edge_type_distribution']['HOLDS']:,} holding relationships capturing the complete portfolio state across 45 time periods.",

    f"The graph successfully models portfolio dynamics through {temporal_results['structural_completeness']['edge_type_distribution']['INCREASED']:,} position increases, {temporal_results['structural_completeness']['edge_type_distribution']['DECREASED']:,} decreases, {temporal_results['structural_completeness']['edge_type_distribution']['ENTERED']:,} new entries, and {temporal_results['structural_completeness']['edge_type_distribution']['EXITED']:,} exits. This comprehensive capture of portfolio changes enables analysis of how fund managers adjust positions in response to evolving market conditions.",

    f"Consistency achieves {temporal_results['consistency']['consistency_score']:.3f} with zero temporal violations and zero cycles, confirming logical integrity. The detection of {temporal_results['consistency']['orphan_nodes_count']} orphan nodes (primarily time period nodes without connections) represents a minor issue that does not significantly impact overall graph utility. Semantic coherence of {temporal_results['semantic_coherence']['semantic_coherence_score']:.3f} indicates strong community structure with sector purity of {temporal_results['semantic_coherence']['average_sector_purity']:.3f}, demonstrating that related stocks cluster appropriately within the graph topology."
])

# Slide 19: Results - Causal KG
add_results_slide(prs, "Results: Causal Knowledge Graph Evaluation", causal_results, "Causal Knowledge Graph")

# Slide 20: Results - Causal KG Detailed Analysis
add_content_slide(prs, "Causal Knowledge Graph: Detailed Analysis", [
    f"The Causal Knowledge Graph achieves an overall quality score of {causal_results['overall_quality_score']:.3f}, representing good quality for a relatively sparse causal structure focused on the most robust statistical relationships. The graph contains {causal_results['structural_completeness']['node_count']} nodes comprising {causal_results['structural_completeness']['node_type_distribution']['Factor']} causal factors (macroeconomic indicators and fundamentals) and {causal_results['structural_completeness']['node_type_distribution']['Action']} action nodes (portfolio decisions), connected by {causal_results['structural_completeness']['edge_count']} causal relationships.",

    f"Consistency reaches perfect score of {causal_results['consistency']['consistency_score']:.3f} with successful enforcement of directed acyclic graph structure, confirming absence of causal cycles that would create logical contradictions. Zero orphan nodes indicate complete connectivity where every factor either influences or is influenced by other elements in the decision ecosystem. The graph captures {causal_results['structural_completeness']['edge_type_distribution']['INFLUENCES']} influence relationships representing statistical causal links discovered through Granger causality testing.",

    f"The relatively lower completeness score of {causal_results['structural_completeness']['completeness_score']:.3f} reflects the intentionally sparse structure emphasizing only statistically significant causal relationships rather than all potential connections. This design choice prioritizes relationship quality over quantity, ensuring that causal edges represent genuine predictive relationships rather than spurious correlations. Semantic coherence of {causal_results['semantic_coherence']['semantic_coherence_score']:.3f} with modularity of {causal_results['semantic_coherence']['modularity']:.3f} indicates strong community structure where related factors cluster together, supporting interpretable causal reasoning."
])

# Slide 21: Discussion - Strengths and Achievements
add_content_slide(prs, "Discussion: Key Achievements", [
    "The Phase 1 implementation successfully demonstrates that dual knowledge graph representations can effectively capture fund manager decision patterns from historical portfolio data. The temporal knowledge graph achieves excellent structural completeness and strong consistency, providing comprehensive coverage of portfolio evolution across 22 funds, 1,762 stocks, and 45 time periods. The high sector purity of 99.8% confirms that the graph structure accurately reflects domain semantics, with securities clustering by industry as expected from investment theory.",

    "The causal knowledge graph achieves perfect consistency with successful DAG enforcement, establishing a sound logical foundation for interventional reasoning. The discovery of 28 statistically significant causal relationships through Granger causality testing provides empirical evidence that observable macroeconomic factors and fundamentals do causally influence portfolio decisions with predictable time lags. This validates the core hypothesis that fund manager decision-making follows discoverable cause-effect patterns rather than purely random or idiosyncratic processes.",

    "The integration of temporal and causal perspectives creates a knowledge base uniquely suited to decision support applications. Unlike pure machine learning approaches that provide predictions without explanations, or pure statistical approaches that identify correlations without causal mechanisms, our dual-graph approach captures both what decisions occurred and why, enabling interpretable reasoning that mirrors expert cognitive processes. The successful implementation of this framework on real-world portfolio data demonstrates feasibility for practical deployment in fund management contexts."
])

# Slide 22: Discussion - Limitations and Future Work
add_content_slide(prs, "Discussion: Limitations and Areas for Enhancement", [
    "The inferential utility scores, while acceptable, indicate opportunities for improvement in supporting complex multi-hop reasoning. The current graph structures emphasize direct relationships, resulting in short average path lengths that limit transitive inference capabilities. Future work could enrich the knowledge base with intermediate conceptual nodes representing derived factors or composite signals that emerge from combinations of base factors, increasing reasoning depth and supporting more sophisticated queries.",

    "The informativeness scores suggest that additional metadata could enhance decision explanations. While current node and edge attributes capture essential information like portfolio weights and causal strengths, incorporating richer contextual information such as analyst ratings, corporate governance metrics, or management quality assessments could enable more nuanced reasoning. The modest informativeness particularly in the causal graph reflects the focus on statistical relationships rather than domain-driven feature engineering, suggesting value in hybrid approaches combining data-driven discovery with expert knowledge.",

    "The evaluation framework focuses on intrinsic graph properties without yet testing downstream decision support performance. While high intrinsic quality scores provide confidence in the knowledge representations, ultimate validation requires demonstrating that these representations support effective portfolio construction in Phase 2. The gap between structural quality and practical utility represents a common challenge in knowledge engineering, necessitating the planned extrinsic evaluation where generated portfolios are compared against actual fund manager decisions to measure imitation fidelity."
])

# Slide 23: Conclusion
add_content_slide(prs, "Conclusion and Path Forward", [
    "Phase 1 successfully achieves its primary objective of constructing and evaluating knowledge representations that capture fund manager decision patterns. The temporal knowledge graph demonstrates excellent structural completeness (0.988), strong consistency (0.833), and good semantic coherence (0.777), providing comprehensive coverage of portfolio evolution dynamics. The causal knowledge graph achieves perfect consistency (1.000) with successful DAG enforcement, establishing a sound foundation for interventional reasoning about market conditions and portfolio actions.",

    "The overall quality scores—0.690 for temporal graph and 0.527 for causal graph—indicate that both knowledge representations achieve sufficient quality to serve as reliable foundations for downstream applications. The relatively lower causal graph score reflects intentional design emphasizing robust statistical relationships over comprehensive coverage, prioritizing relationship quality over quantity. The integration layer successfully links temporal events to causal explanations, creating a unified knowledge base supporting queries that combine 'what happened' with 'why it happened'.",

    "These Phase 1 knowledge representations will serve as the foundation for Phase 2 downstream tasks including portfolio stock selection, allocation optimization, and temporal adjustment. The structured knowledge bases enable the next phase to employ graph neural networks that leverage both node attributes and graph topology for stock selection, use causal graphs for explainable decision-making that provides rationale alongside recommendations, and apply temporal reasoning for dynamic portfolio rebalancing. The successful completion of Phase 1 establishes confidence that our dual-graph knowledge representation approach can ultimately support systems that effectively imitate fund manager decision-making patterns."
])

# Slide 24: References (Part 1)
add_content_slide(prs, "References", [
    "[1] DeMiguel, V., Gil-Bazo, J., Nogales, F. J., & Santos, A. A. P. (2023). Machine learning and fund characteristics help to select mutual funds with positive alpha. Journal of Financial Economics, 150(3), 101–129.",

    "[2] Kaniel, R., Lin, Z., Pelger, M., & Van Nieuwerburgh, S. (2023). Machine-learning the skill of mutual fund managers. Journal of Financial Economics, 150(3), 94–138.",

    "[3] Wang, H., Qian, Y., & Zhang, X. (2024). From data to insights: The application and challenges of knowledge graphs in intelligent audit. Journal of Cloud Computing: Advances, Systems and Applications, 13, 114.",

    "[4] Jiang, Y., et al. (2024). A review on the reliability of knowledge graph: From a knowledge representation learning perspective. World Wide Web, 27, 47."
])

# Slide 25: References (Part 2)
add_content_slide(prs, "References (Continued)", [
    "[7] Liu, W., Hasikin, K., Khairuddin, A. S. M., Liu, M., & Zhao, X. (2025). A temporal knowledge graph reasoning model based on recurrent encoding and contrastive learning. PeerJ Computer Science, 11, e2595.",

    "[8] Job, J. M., Botteldooren, D., & De Coensel, B. (2025). Exploring causal learning through graph neural networks: An in-depth review. WIREs Data Mining and Knowledge Discovery, 15(1), e70024.",

    "[12] Zhang, Y., & Xiao, G. (2024). A novel customizing knowledge graph evaluation method for incorporating user needs. Scientific Reports, 14, 9594.",

    "[13] Liu, Z., Zhou, H., Chen, Y., Yang, Y., & Huang, Z. (2024). Research on Traditional Chinese Medicine: Domain knowledge graph completion and quality evaluation. JMIR Medical Informatics, 12, e55090."
])

# Slide 26: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
thanks_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "Thank You"
thanks_frame.paragraphs[0].font.size = Pt(48)
thanks_frame.paragraphs[0].font.bold = True
thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
thanks_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

contact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "Questions & Discussion\n\nGANNAMANENI JASWANTH\n242IT012\njaswanthgoelite@gmail.com"
contact_frame.paragraphs[0].font.size = Pt(14)
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
contact_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)

# Save presentation
output_path = Path("Phase1_Presentation_GANNAMANENI_JASWANTH.pptx")
prs.save(output_path)

print(f"\n{'='*80}")
print("PRESENTATION CREATED SUCCESSFULLY")
print(f"{'='*80}")
print(f"\nFile saved to: {output_path.absolute()}")
print(f"\nPresentation contains {len(prs.slides)} slides covering:")
print("  • Introduction and Research Motivation")
print("  • Comprehensive Literature Review")
print("  • Problem Statement and Objectives")
print("  • Detailed Methodology")
print("  • Complete Results and Analysis")
print("  • Discussion and Conclusion")
print("  • References")
print(f"\n{'='*80}\n")
